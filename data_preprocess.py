import os
import subprocess
from pptx import Presentation
from docx import Document
import fitz  # PyMuPDF


def rename_files_with_numbers(folder_path):
    if not os.path.isdir(folder_path):
        print(f"❌ 文件夹不存在: {folder_path}")
        return

    allowed_exts = {'.ppt', '.pptx', '.doc', '.docx', '.pdf'}
    files = sorted(os.listdir(folder_path))
    count = 1

    for filename in files:
        old_path = os.path.join(folder_path, filename)
        if not os.path.isfile(old_path):
            continue

        ext = os.path.splitext(filename)[1].lower()
        if ext not in allowed_exts:
            continue

        if ext == '.pdf':
            new_filename = f"pdf_{count}_{filename}"
        else:
            new_filename = f"{count}_{filename}"

        new_path = os.path.join(folder_path, new_filename)
        os.rename(old_path, new_path)
        print(f"✅ 重命名: {filename} -> {new_filename}")
        count += 1


def convert_with_libreoffice(input_file, libreoffice_path, docx_output_folder):
    if not os.path.exists(input_file):
        print(f"❌ 文件不存在: {input_file}")
        return False

    ext = os.path.splitext(input_file)[1].lower()
    if ext == ".ppt":
        target_ext = "pptx"
        output_dir = os.path.dirname(input_file)
    elif ext == ".doc":
        target_ext = "docx"
        output_dir = docx_output_folder
    else:
        print(f"⚠️ 不支持的文件类型: {input_file}")
        return False

    try:
        subprocess.run([
            libreoffice_path,
            "--headless",
            "--convert-to", target_ext,
            input_file,
            "--outdir", output_dir
        ], check=True)

        print(f"✅ 成功转换: {input_file} → .{target_ext}")

        # 获取转换后文件路径
        base_name = os.path.splitext(os.path.basename(input_file))[0]
        converted_file = os.path.join(output_dir, f"{base_name}.{target_ext}")

        # 检查文件是否成功生成
        if os.path.exists(converted_file):
            os.remove(input_file)
            print(f"🗑️ 已删除原文件: {input_file}")
            return True
        else:
            print(f"⚠️ 转换后文件未找到，未删除原文件: {converted_file}")
            return False

    except subprocess.CalledProcessError as e:
        print(f"❌ 转换失败: {input_file}\n错误信息: {e}")
        return False

def batch_convert_folder(input_folder, libreoffice_path, docx_output_folder):
    if not os.path.exists(input_folder):
        print(f"❌ 输入文件夹不存在: {input_folder}")
        return

    os.makedirs(docx_output_folder, exist_ok=True)

    for filename in os.listdir(input_folder):
        filepath = os.path.join(input_folder, filename)
        if os.path.isfile(filepath) and filename.lower().endswith((".ppt", ".doc")):
            convert_with_libreoffice(filepath, libreoffice_path, docx_output_folder)


def extract_text_from_ppt(ppt_path):
    prs = Presentation(ppt_path)
    text = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)

    return "\n".join(text)

def convert_ppt_to_docx(input_folder, output_folder):
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    for filename in os.listdir(input_folder):
        if filename.lower().endswith(('.ppt', '.pptx')):
            ppt_path = os.path.join(input_folder, filename)
            content = extract_text_from_ppt(ppt_path)

            # 创建 Word 文档并写入内容
            doc = Document()
            for line in content.splitlines():
                if line.strip():
                    doc.add_paragraph(line.strip())

            # 保存为 .docx 文件
            docx_filename = os.path.splitext(filename)[0] + ".docx"
            output_path = os.path.join(output_folder, docx_filename)
            doc.save(output_path)
            print(f"✅ 已保存: {output_path}")

            # 删除原始文件
            try:
                os.remove(ppt_path)
                print(f"🗑️ 已删除原文件: {ppt_path}")
            except Exception as e:
                print(f"⚠️ 删除失败: {ppt_path}\n错误信息: {e}")



def extract_text_from_pdf(pdf_path, handle_columns=True):
    doc = fitz.open(pdf_path)
    full_text = []

    for page_num in range(len(doc)):
        page = doc[page_num]

        if handle_columns:
            blocks = page.get_text("blocks")
            blocks.sort(key=lambda b: (b[1], b[0]))  # 按 y（行）和 x（列）排序
            page_text = [block[4].strip() for block in blocks if block[4].strip()]
            full_text.extend(page_text)
        else:
            page_text = page.get_text().split('\n')
            full_text.extend(line.strip() for line in page_text if line.strip())

    return "\n".join(full_text)

def save_text_to_docx(text, output_path):
    doc = Document()
    for para in text.split("\n"):
        if para.strip():
            doc.add_paragraph(para.strip())
    doc.save(output_path)

def batch_convert_pdf_to_docx(input_folder, output_folder, handle_columns=True):
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    for filename in os.listdir(input_folder):
        if filename.lower().endswith(".pdf"):
            pdf_path = os.path.join(input_folder, filename)
            docx_name = os.path.splitext(filename)[0] + ".docx"
            output_path = os.path.join(output_folder, docx_name)

            print(f"📄 正在处理: {filename}")
            try:
                text = extract_text_from_pdf(pdf_path, handle_columns=handle_columns)
                save_text_to_docx(text, output_path)
                os.remove(pdf_path)  # ✅ 删除原始 PDF 文件
                print(f"✅ 已保存为: {output_path}，并删除原文件")
            except Exception as e:
                print(f"❌ 处理失败 {filename}，错误：{e}")


import os
import fitz  # PyMuPDF
from docx import Document

def extract_two_column_text(page):
    """
    提取一页中左栏和右栏的文本，左栏在前，右栏在后
    """
    width = page.rect.width
    height = page.rect.height

    mid_x = width / 2  # 中线分为两栏
    left_rect = fitz.Rect(0, 0, mid_x, height)
    right_rect = fitz.Rect(mid_x, 0, width, height)

    # 提取左栏和右栏文字
    left_text = page.get_textbox(left_rect)
    right_text = page.get_textbox(right_rect)

    return left_text.strip() + '\n' + right_text.strip()

def extract_text_from_pdf2(pdf_path):
    doc = fitz.open(pdf_path)
    full_text = []

    for page in doc:
        page_text = extract_two_column_text(page)
        full_text.append(page_text)

    return '\n\n'.join(full_text)

def convert_pdfs_to_docx(input_folder, output_folder):
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)

    for filename in os.listdir(input_folder):
        if filename.lower().endswith('.pdf'):
            pdf_path = os.path.join(input_folder, filename)
            text = extract_text_from_pdf2(pdf_path)

            doc = Document()
            for para in text.split('\n'):
                if para.strip():
                    doc.add_paragraph(para.strip())

            docx_filename = os.path.splitext(filename)[0] + '.docx'
            docx_path = os.path.join(output_folder, docx_filename)
            doc.save(docx_path)

            # 删除原始 PDF 文件
            os.remove(pdf_path)
            print(f"✅ 转换并删除原文件：{docx_path}")



def data_preprocessing():

    libreoffice_path = r"D:\software\LibreOffice\program\soffice.exe"
    input_folder = r"D:\projects\sailisi\图谱库语料"
    output_folder = r"D:\projects\sailisi\切片"
    two_column_pdf = r"D:\projects\sailisi\两栏pdf" #这里放置两栏pdf文件


    rename_files_with_numbers(input_folder)

    rename_files_with_numbers(two_column_pdf)

    # 1.将文件夹中的所有doc和ppt格式全部转化为docx和pptx格式。并删除原先文件
    batch_convert_folder(input_folder, libreoffice_path, output_folder)
 
    # 2. 将文件夹中的ppt全部转化为docx文档，保存在切片文件夹中，并删除原先文件
    convert_ppt_to_docx(input_folder, output_folder)

    # 3. 将文件夹中的pdf（不包含双栏pdf）内容转化为docx文档，并删除原先文件
    batch_convert_pdf_to_docx(input_folder, output_folder, handle_columns=True)

    # 4.将文件夹中的pdf(只是两栏pdf)内容转化为docx文档，该操作不删除原先文件
    convert_pdfs_to_docx(two_column_pdf, output_folder)
    
if __name__ =="__main__":
    data_preprocessing()